import os
import tempfile
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def pdf_to_pptx(pdf_file):
    # Save the PDF file to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf_file:
        temp_pdf_file.write(pdf_file.read())

    
    # Get the path of the saved PDF file
    pdf_file_path = os.path.abspath(temp_pdf_file.name)

    # Convert the PDF file to images using pdf2image
    images = convert_from_path(pdf_file_path)

    # Create a new PowerPoint presentation
    prs = Presentation()

    # Add each image as a new slide in the presentation
    for image in images:
        # Save the image to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_image_file:
            image.save(temp_image_file.name, format='PNG')

        # Get the path of the saved image file
        image_file_path = os.path.abspath(temp_image_file.name)

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add the image to the slide
        pic = slide.shapes.add_picture(image_file_path, 0, 0, prs.slide_width, prs.slide_height)
        pic.crop_left = Inches(0)
        pic.crop_right = Inches(0)
        pic.crop_top = Inches(0)
        pic.crop_bottom = Inches(0)
        pic.width = prs.slide_width
        pic.height = prs.slide_height

        # Delete the temporary image file
        os.unlink(image_file_path)

    # Save the PowerPoint presentation to a new file
    pptx_file_path = os.path.splitext(pdf_file_path)[0] + '.pptx'
    prs.save(pptx_file_path)

    # Delete the temporary PDF file
    os.unlink(pdf_file_path)

    # Read the PPTX file content and return it
    with open(pptx_file_path, 'rb') as pptx_file:
        pptx_content = pptx_file.read()

    # Delete the temporary PPTX file
    os.unlink(pptx_file_path)

    return pptx_content